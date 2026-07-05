import os
import io
import re
import pandas as pd
from typing import List, Dict, Any
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_groq import ChatGroq
from langchain_core.documents import Document as LangchainDocument
from langchain_core.prompts import PromptTemplate

from mcp import MCPMessage


class IngestionAgent:
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1200,
            chunk_overlap=200,
            separators=["\n\n", "\n", ". ", " ", ""]
        )

    def _clean_text(self, text: str) -> str:
        """Clean messy text while PRESERVING line structure (bullets, numbering, headings, spacing)."""
        # Remove control characters (but keep \n and \t)
        text = re.sub(r'[\x00-\x08\x0b-\x0c\x0e-\x1f\x7f]', '', text)
        # Strip URLs and emails
        text = re.sub(r'http\S+|www\.\S+', '', text)
        text = re.sub(r'\S+@\S+', '', text)

        cleaned_lines = []
        for line in text.split('\n'):
            # Collapse only intra-line whitespace (spaces/tabs), never the newline itself
            line = re.sub(r'[ \t]+', ' ', line)
            # Collapse repeated punctuation like "...." or "??" -> "."
            line = re.sub(r'[.!?]{2,}', '.', line)
            cleaned_lines.append(line.strip())

        text = '\n'.join(cleaned_lines)
        # Collapse 3+ consecutive blank lines down to a single blank line (paragraph break)
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

    def _extract_text(self, file_content: bytes, filename: str) -> str:
        """Extract text with quality control, preserving structural line breaks."""
        ext = filename.split('.')[-1].lower()
        text = ""

        try:
            if ext == "pdf":
                reader = PdfReader(io.BytesIO(file_content))

                if len(reader.pages) == 0:
                    raise ValueError("PDF has no pages")

                for page_num, page in enumerate(reader.pages, 1):
                    try:
                        extracted = page.extract_text()
                        if extracted and len(extracted.strip()) > 10:
                            cleaned = self._clean_text(extracted)
                            if cleaned:
                                text += f"PAGE {page_num}:\n{cleaned}\n\n"
                    except Exception:
                        continue

            elif ext in ["docx", "doc"]:
                doc = Document(io.BytesIO(file_content))
                for para in doc.paragraphs:
                    raw = para.text
                    if raw.strip() and len(raw.strip()) > 1:
                        # Detect list-style paragraphs (Word list styles / numbering)
                        style_name = (para.style.name or "").lower()
                        is_list_item = (
                            "list" in style_name
                            or "bullet" in style_name
                            or bool(re.match(r'^\s*(\d+[\.\)]|[-*•])\s+', raw))
                        )
                        cleaned = self._clean_text(raw)
                        if cleaned:
                            if is_list_item and not re.match(r'^\s*(\d+[\.\)]|[-*•])\s+', cleaned):
                                cleaned = f"- {cleaned}"
                            # Headings get their own paragraph break before/after
                            if "heading" in style_name or "title" in style_name:
                                text += f"\n{cleaned}\n\n"
                            else:
                                text += cleaned + "\n"

            elif ext in ["pptx", "ppt"]:
                ppt = Presentation(io.BytesIO(file_content))
                for slide_num, slide in enumerate(ppt.slides, 1):
                    slide_text = f"SLIDE {slide_num}:\n"
                    for shape in slide.shapes:
                        if not hasattr(shape, "text") or not shape.text.strip():
                            continue
                        # Preserve per-paragraph structure within a text frame (bullets)
                        if hasattr(shape, "text_frame"):
                            for para in shape.text_frame.paragraphs:
                                raw_line = "".join(run.text for run in para.runs) or para.text
                                if raw_line.strip():
                                    cleaned_line = self._clean_text(raw_line)
                                    if cleaned_line:
                                        indent = "  " * (para.level or 0)
                                        slide_text += f"{indent}- {cleaned_line}\n"
                        else:
                            cleaned = self._clean_text(shape.text)
                            if cleaned:
                                slide_text += cleaned + "\n"
                    if len(slide_text) > 15:
                        text += slide_text + "\n"

            elif ext == "csv":
                df = pd.read_csv(io.BytesIO(file_content))
                text = df.to_string()

            elif ext in ["txt", "md"]:
                text = file_content.decode("utf-8", errors="ignore")
                text = self._clean_text(text)

        except Exception as e:
            text = f"Error: {str(e)}"

        return text.strip()

    def _remove_duplicate_content(self, text: str) -> str:
        """Remove repeated lines while preserving order and structure."""
        lines = text.split('\n')
        seen_lines = set()
        unique_lines = []

        for line in lines:
            # Always keep blank lines / very short lines (structure), don't dedupe them
            if len(line.strip()) < 3:
                unique_lines.append(line)
                continue

            normalized = line.strip().lower()

            if normalized not in seen_lines:
                seen_lines.add(normalized)
                unique_lines.append(line)

        return '\n'.join(unique_lines)

    def process(self, message: MCPMessage) -> MCPMessage:
        if message.type != "INGEST_REQUEST":
            raise ValueError(f"IngestionAgent does not support {message.type}")

        file_content = message.payload.get("file_content")
        filename = message.payload.get("filename")

        if not file_content or not filename:
            raise ValueError("Missing file_content or filename")

        raw_text = self._extract_text(file_content, filename)

        if not raw_text or len(raw_text) < 20:
            raise ValueError(f"Could not extract meaningful text from {filename}")

        clean_text = self._remove_duplicate_content(raw_text)

        docs = [LangchainDocument(
            page_content=clean_text,
            metadata={"source": filename, "type": filename.split('.')[-1].lower()}
        )]

        chunks = self.text_splitter.split_documents(docs)

        if not chunks:
            raise ValueError(f"No chunks from {filename}")

        payload_chunks = [{"page_content": c.page_content, "metadata": c.metadata} for c in chunks]

        return MCPMessage(
            sender="IngestionAgent",
            receiver="RetrievalAgent",
            type="DOCUMENT_CHUNKS",
            trace_id=message.trace_id,
            payload={"chunks": payload_chunks, "filename": filename, "chunk_count": len(chunks)}
        )


class RetrievalAgent:
    def __init__(self, embedding_model="sentence-transformers/all-MiniLM-L6-v2"):
        self.embeddings = HuggingFaceEmbeddings(model_name=embedding_model, model_kwargs={"device": "cpu"})
        self.vector_store = None

    def _is_duplicate_chunk(self, chunk1: str, chunk2: str) -> bool:
        """Check for duplicate chunks"""
        norm1 = ' '.join(chunk1.lower().split())
        norm2 = ' '.join(chunk2.lower().split())

        if norm1 == norm2:
            return True

        if len(norm1) > 50 and len(norm2) > 50:
            if norm1[:100] == norm2[:100]:
                return True

        return False

    def _deduplicate_results(self, results: List) -> List:
        """Remove duplicates"""
        deduped = []

        for result in results:
            is_duplicate = False

            for existing in deduped:
                if self._is_duplicate_chunk(result.page_content, existing.page_content):
                    is_duplicate = True
                    break

            if not is_duplicate:
                deduped.append(result)

        return deduped

    def process(self, message: MCPMessage) -> MCPMessage:
        if message.type == "DOCUMENT_CHUNKS":
            chunks_data = message.payload.get("chunks", [])

            if not chunks_data:
                return MCPMessage(
                    sender="RetrievalAgent",
                    receiver="System",
                    type="INGESTION_COMPLETE",
                    trace_id=message.trace_id,
                    payload={"status": "error", "chunks_added": 0}
                )

            docs = [LangchainDocument(page_content=c["page_content"], metadata=c.get("metadata", {})) for c in chunks_data]

            try:
                if self.vector_store is None:
                    self.vector_store = FAISS.from_documents(docs, self.embeddings)
                else:
                    self.vector_store.add_documents(docs)

                return MCPMessage(
                    sender="RetrievalAgent",
                    receiver="System",
                    type="INGESTION_COMPLETE",
                    trace_id=message.trace_id,
                    payload={"status": "success", "chunks_added": len(docs)}
                )
            except Exception as e:
                return MCPMessage(
                    sender="RetrievalAgent",
                    receiver="System",
                    type="INGESTION_COMPLETE",
                    trace_id=message.trace_id,
                    payload={"status": "error", "chunks_added": 0, "error": str(e)}
                )

        elif message.type == "RETRIEVE_REQUEST":
            query = message.payload.get("query", "")
            top_k = message.payload.get("top_k", 5)

            if not query:
                return MCPMessage(
                    sender="RetrievalAgent",
                    receiver="LLMResponseAgent",
                    type="CONTEXT_RESPONSE",
                    trace_id=message.trace_id,
                    payload={"query": query, "context_chunks": [], "error": "Empty query"}
                )

            if not self.vector_store:
                return MCPMessage(
                    sender="RetrievalAgent",
                    receiver="LLMResponseAgent",
                    type="CONTEXT_RESPONSE",
                    trace_id=message.trace_id,
                    payload={"query": query, "context_chunks": [], "error": "No documents"}
                )

            try:
                results = self.vector_store.similarity_search(query, k=top_k + 3)
                unique_results = self._deduplicate_results(results)
                unique_results = unique_results[:top_k]

                context_chunks = [
                    {
                        "page_content": doc.page_content,
                        "metadata": doc.metadata
                    }
                    for doc in unique_results
                ]

                return MCPMessage(
                    sender="RetrievalAgent",
                    receiver="LLMResponseAgent",
                    type="CONTEXT_RESPONSE",
                    trace_id=message.trace_id,
                    payload={"query": query, "context_chunks": context_chunks}
                )
            except Exception as e:
                return MCPMessage(
                    sender="RetrievalAgent",
                    receiver="LLMResponseAgent",
                    type="CONTEXT_RESPONSE",
                    trace_id=message.trace_id,
                    payload={"query": query, "context_chunks": [], "error": str(e)}
                )
        else:
            raise ValueError(f"RetrievalAgent does not support {message.type}")


class LLMResponseAgent:
    def __init__(self, groq_api_key: str, model_name: str = "llama-3.3-70b-versatile"):
        if not groq_api_key:
            raise ValueError("Groq API key required")

        self.llm = ChatGroq(
            temperature=0.0,  # ZERO CREATIVITY - facts only
            model_name=model_name,
            groq_api_key=groq_api_key,
            max_tokens=1000,
            timeout=30
        )

        # STRICT PROMPT - facts only, but formatting-preserving
        self.prompt_template = PromptTemplate(
            input_variables=["context", "query"],
            template="""EXTRACT ONLY THE FACTS FROM THE DOCUMENT.

FORBIDDEN:
- Do NOT add "This indicates that"
- Do NOT add "However"
- Do NOT add "Unfortunately"
- Do NOT add "In summary"
- Do NOT add analysis or commentary
- Do NOT say what information is missing
- Do NOT mention other documents
- Do NOT provide explanations
- Do NOT add any interpretation
- Do NOT merge bullet points, numbered lists, or headings into a single paragraph
- Do NOT remove line breaks that separate list items or sections

REQUIRED:
- Only state the facts from the document
- Only answer the question asked
- Preserve the ORIGINAL FORMATTING exactly as it appears in the DOCUMENT TEXT below
- If the source uses bullet points, reproduce them as bullet points (one per line, starting with "-")
- If the source uses numbered lists, reproduce them as numbered lists (1., 2., 3., each on its own line)
- If the source uses headings, keep them on their own line, followed by a blank line
- Do not paraphrase
- Do not summarize
- Do not analyze

DOCUMENT TEXT:
{context}

USER QUESTION:
{query}

FACTS ONLY, WITH ORIGINAL FORMATTING PRESERVED (bullets/numbering/headings/line breaks intact):"""
        )

    def _normalize_markdown(self, text: str) -> str:
        """
        Force proper Markdown list formatting so bullets/numbered lists render
        vertically instead of collapsing into one horizontal paragraph.

        Handles two common failure modes from LLM output:
        1. A list item follows a normal sentence with no blank line before it
           -> Markdown treats it as part of the same paragraph.
        2. Multiple list items end up on the same physical line
           (e.g. "1. foo 2. bar 3. baz") -> split them onto separate lines.
        """
        # Split "1. foo 2. bar" into separate lines when numbers appear mid-line
        text = re.sub(r'(?<!^)(?<!\n)\s(?=\d+\.\s)', '\n', text)
        # Split " - foo - bar" into separate lines when dashes appear mid-line
        text = re.sub(r'(?<!^)(?<!\n)\s(?=-\s\S)', '\n', text)

        lines = text.split('\n')
        normalized_lines = []
        list_marker = re.compile(r'^\s*(\d+[\.\)]|[-*•])\s+')

        for i, line in enumerate(lines):
            is_list_item = bool(list_marker.match(line))
            prev_line = normalized_lines[-1] if normalized_lines else ""
            prev_is_list_item = bool(list_marker.match(prev_line))

            # Insert a blank line right before a list starts (but not between
            # consecutive list items, and not at the very start of the text)
            if is_list_item and normalized_lines and prev_line.strip() and not prev_is_list_item:
                normalized_lines.append("")

            normalized_lines.append(line)

        result = '\n'.join(normalized_lines)
        # Collapse any accidental triple+ blank lines back down
        result = re.sub(r'\n{3,}', '\n\n', result)
        return result.strip()

    def process(self, message: MCPMessage) -> MCPMessage:
        if message.type != "CONTEXT_RESPONSE":
            raise ValueError(f"LLMResponseAgent does not support {message.type}")

        query = message.payload.get("query", "")
        context_chunks = message.payload.get("context_chunks", [])
        error = message.payload.get("error")

        if error or not context_chunks:
            error_msg = error or "No information found"
            return MCPMessage(
                sender="LLMResponseAgent",
                receiver="UI",
                type="FINAL_ANSWER",
                trace_id=message.trace_id,
                payload={"answer": error_msg, "sources": [], "status": "error"}
            )

        try:
            # Build context - ONLY relevant chunks, full content (no mid-list truncation)
            context_parts = []
            for chunk in context_chunks:
                source = chunk['metadata'].get('source', 'Unknown')
                content = chunk['page_content'].strip()
                context_parts.append(f"[{source}]\n{content}")

            context_text = "\n\n".join(context_parts)

            # Raised limit so lists/headings aren't cut mid-item; trim on a line boundary
            MAX_CONTEXT_CHARS = 6000
            if len(context_text) > MAX_CONTEXT_CHARS:
                truncated = context_text[:MAX_CONTEXT_CHARS]
                last_newline = truncated.rfind('\n')
                if last_newline > 0:
                    truncated = truncated[:last_newline]
                context_text = truncated

            prompt = self.prompt_template.format(context=context_text, query=query)

            response = self.llm.invoke(prompt)
            response_text = response.content.strip()

            if not response_text:
                response_text = "No information found"
            else:
                response_text = self._normalize_markdown(response_text)

            return MCPMessage(
                sender="LLMResponseAgent",
                receiver="UI",
                type="FINAL_ANSWER",
                trace_id=message.trace_id,
                payload={"answer": response_text, "sources": context_chunks, "status": "success"}
            )

        except Exception as e:
            return MCPMessage(
                sender="LLMResponseAgent",
                receiver="UI",
                type="FINAL_ANSWER",
                trace_id=message.trace_id,
                payload={"answer": f"Error: {str(e)[:80]}", "sources": [], "status": "error"}
            )
